"""
Build the Week 1 PPTX presentation (Exploits University, BICT 3202 OOAD).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from common import BRAND, VISUALS, OUTPUT

MAROON = RGBColor(0x7B, 0x11, 0x13)
DARK   = RGBColor(0x59, 0x09, 0x0C)
SOFT   = RGBColor(0xA9, 0x44, 0x42)
GOLD   = RGBColor(0xC9, 0xA2, 0x27)
GOLD_L = RGBColor(0xE4, 0xC9, 0x7C)
CREAM  = RGBColor(0xFA, 0xF6, 0xEC)
CREAM2 = RGBColor(0xF3, 0xEC, 0xDC)
TINT   = RGBColor(0xF3, 0xE4, 0xE4)
INK    = RGBColor(0x2B, 0x26, 0x26)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

H_FONT = "Poppins"
B_FONT = "Open Sans"

SW = 13.333
SH = 7.5

_page = [0]


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, runs, size=18, color=INK, font=B_FONT, bold=False,
          italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          line_spacing=1.0, space_after=0):
    """runs: str  OR  list of paragraphs; each paragraph is str or list of (txt, overrides)"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, ov in para:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.color.rgb = ov.get("color", color)
    return box


def _pic(slide, name, x=None, y=None, w=None):
    path = os.path.join(VISUALS, name)
    if x is None:
        x = (SW - w) / 2.0
    kw = dict(left=Inches(x), top=Inches(y), width=Inches(w))
    return slide.shapes.add_picture(path, **kw)


def _header(slide, title, subtitle=None):
    _rect(slide, 0, 0, SW, 1.18, MAROON)
    _rect(slide, 0, 1.18, SW, 0.06, GOLD)
    _rect(slide, 0, 1.24, SW, 0.02, DARK)
    _text(slide, 0.55, 0.16, 10.6, 0.6, title, size=27, color=WHITE,
          font=H_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _text(slide, 0.55, 0.72, 10.6, 0.4, subtitle, size=13, color=GOLD_L,
              font=B_FONT, anchor=MSO_ANCHOR.MIDDLE)
    # footer
    _page[0] += 1
    _text(slide, 0.55, 7.12, 8, 0.3,
          "BICT 3202 · Object-Oriented Analysis and Design", size=10, color=GREY)
    _text(slide, 9.5, 7.12, 3.3, 0.3, f"Week 1 · Slide {_page[0]}", size=10,
          color=GREY, align=PP_ALIGN.RIGHT)


def _takeaway(slide, text):
    _rect(slide, 0.9, 6.55, SW - 1.8, 0.5, DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, 1.2, 6.55, SW - 2.4, 0.5, text, size=14, color=GOLD_L,
          font=B_FONT, bold=False, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def _bullets(slide, x, y, w, h, items, size=16, gap=8, color=INK):
    """items: list of str or (lead, rest); lead rendered bold maroon."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.06
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "▪  "
        r.font.name = B_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = MAROON
        r.font.bold = True
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead
            r1.font.name = B_FONT; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = MAROON
            r2 = p.add_run(); r2.text = rest
            r2.font.name = B_FONT; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r1 = p.add_run(); r1.text = it
            r1.font.name = B_FONT; r1.font.size = Pt(size)
            r1.font.color.rgb = color
    return box


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]
    seal = os.path.join(BRAND, "eu_seal.png")

    # ---------------- 1. COVER ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, MAROON)
    _rect(s, 0, 0, SW, 7.5, MAROON)
    _rect(s, 0.35, 0.35, SW - 0.7, 6.8, MAROON, line=GOLD)
    _rect(s, 0, 6.1, SW, 1.4, DARK)
    s.shapes.add_picture(seal, Inches((SW - 1.7) / 2), Inches(0.75), Inches(1.7), Inches(1.7))
    _text(s, 0.5, 2.6, SW - 1, 0.6, "EXPLOITS UNIVERSITY", size=34, color=GOLD_L,
          font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
    _text(s, 0.5, 3.3, SW - 1, 0.4, "BSc INFORMATION COMMUNICATION AND TECHNOLOGY",
          size=14, color=CREAM, font=B_FONT, align=PP_ALIGN.CENTER)
    _text(s, 0.5, 3.85, SW - 1, 0.4, "BICT 3202 · OBJECT-ORIENTED ANALYSIS AND DESIGN",
          size=16, color=GOLD, font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
    _text(s, 0.5, 4.5, SW - 1, 1.0, "WEEK 1", size=54, color=WHITE, font=H_FONT,
          bold=True, align=PP_ALIGN.CENTER)
    _text(s, 0.5, 5.35, SW - 1, 0.5, "Business Needs, Software Needs and Object Technology for Business",
          size=17, color=CREAM, font=H_FONT, align=PP_ALIGN.CENTER)
    _text(s, 0.5, 6.35, SW - 1, 0.4, "Francis Fweta  ·  ICT Lecturer  ·  Week 1 of 16",
          size=13, color=GOLD_L, font=B_FONT, align=PP_ALIGN.CENTER)

    # ---------------- 2. LEARNING OUTCOMES ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Learning Outcomes", "By the end of this week, you should be able to…")
    _bullets(s, 0.9, 1.6, 11.6, 5.2, [
        ("Define a business need ", "and explain why organisations build or acquire software."),
        ("Distinguish ", "business problem · business need · requirement · software solution."),
        ("Identify stakeholders ", "and explain functional vs non-functional software needs."),
        ("Relate business objectives ", "to software requirements through the business-to-software chain."),
        ("Describe object technology ", "and explain what an object is."),
        ("Explain ", "identity, state and behaviour — and why real-world entities help in analysis."),
        ("Discuss strengths & limitations ", "of object-oriented approaches."),
        ("Apply the concepts ", "to a real-world business case (e.g. registration, ATM, library)."),
    ], size=17, gap=14)
    _takeaway(s, "Requirements come first — not the class diagram.")

    # ---------------- 3. WHY OOAD ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Why Do We Need OOAD?", "The story of a developer who started coding too early")
    _bullets(s, 0.9, 1.55, 5.4, 4.8, [
        "University: “We need a new student management system.”",
        "Developer builds login, tables, payments, exams, dashboard.",
        "Six months later: “This isn't what we asked for.”",
        ("The real problem? ", "insufficient analysis of business needs before design."),
        ("OOAD exists ", "to prevent exactly this failure."),
    ], size=15, gap=12)
    _pic(s, "chart_failure_reasons.png", x=6.5, y=1.5, w=6.4)
    _takeaway(s, "Analysis first — programming second.")

    # ---------------- 4. JOURNEY ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "The Business-to-Software Journey", "Eight steps from a goal to working software")
    _pic(s, "journey.png", y=1.55, w=10.6)
    _takeaway(s, "This chain is the backbone of the entire 16-week module.")

    # ---------------- 5. PROBLEM / NEED / SOLUTION ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Problem → Need → Solution", "Three ideas that must never be confused")
    _pic(s, "problem_need_solution.png", y=1.6, w=10.8)
    _takeaway(s, "Problem = what is wrong · Need = what must improve · Solution = how we improve it")

    # ---------------- 6. BUSINESS OBJECTIVES & CHAIN ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Business Objectives & the Chain", "Objectives drive needs, which drive requirements")
    _bullets(s, 0.9, 1.55, 6.0, 4.8, [
        ("Business objective: ", "a specific result the organisation wants (e.g. cut account-opening from 45 to 10 minutes)."),
        ("Business need: ", "what must improve to reach the objective."),
        ("Not every problem ", "is a software problem — investigate first."),
        ("Constraints & business rules ", "shape the eventual solution."),
    ], size=16, gap=14)
    _bullets(s, 7.1, 1.55, 5.4, 4.8, [
        "Business Objective", "Business Problem", "Business Need", "Stakeholder Needs",
        "Software Requirements", "Analysis Models", "Design", "Implementation",
    ], size=16, gap=8)
    _takeaway(s, "Objective → Problem → Need → Requirements → Models → Design → Implementation")

    # ---------------- 7. STAKEHOLDERS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Stakeholders", "Who has an interest in, or is affected by, the system?")
    _pic(s, "stakeholders.png", y=1.6, w=10.0)
    _takeaway(s, "Ignore a stakeholder and you miss their requirements.")

    # ---------------- 8. ASPECTS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Aspects of Business Software Needs", "A complete need has four parts")
    _pic(s, "aspects_overview.png", y=1.7, w=11.0)
    _takeaway(s, "Functional · Non-functional · Business rules · Constraints")

    # ---------------- 9. FUNCTIONAL vs NON-FUNCTIONAL ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Functional vs Non-Functional", "What it does vs how well it does it")
    _pic(s, "functional_nonfunctional.png", y=1.5, w=10.2)
    _takeaway(s, "Functional = behaviour · Non-functional = quality & constraints")

    # ---------------- 10. BUSINESS RULES & CONSTRAINTS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Business Rules & Constraints", "Policies the software must enforce, and limits on the solution")
    _bullets(s, 0.9, 1.55, 5.9, 4.9, [
        ("Business rule — ", "a policy or condition governing the organisation."),
        "“Pass BICT 2101 before registering for BICT 3202.”",
        "“Withdraw at most MWK 500,000 per day.”",
        "“An employee cannot approve their own claim.”",
        ("The functional requirement ", "tells the software how to enforce the rule."),
    ], size=16, gap=12)
    _bullets(s, 7.1, 1.55, 5.3, 4.9, [
        ("Constraints ", "restrict the solution:"),
        "Available budget", "Existing hardware & database",
        "Legal & organisational policy", "Required platform & deployment",
        "Internet availability",
    ], size=16, gap=10)
    _takeaway(s, "Rules come from the organisation — the software enforces them.")

    # ---------------- 11. ELICITATION ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Requirements Elicitation", "How the analyst discovers stakeholder needs (ISO/IEC/IEEE 29148)")
    _pic(s, "elicitation.png", y=1.65, w=10.8)
    _takeaway(s, "Interviews · Questionnaires · Observation · Documents · Workshops · Prototyping")

    # ---------------- 12. OBJECTS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "What Is an Object?", "Identity · State · Behaviour")
    _pic(s, "object_anatomy.png", y=1.5, w=10.2)
    _takeaway(s, "State = data · Behaviour = methods · Identity = uniqueness")

    # ---------------- 13. CLASS vs OBJECT ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Class vs Object", "The blueprint and the things built from it")
    _pic(s, "class_vs_object.png", y=1.5, w=10.2)
    _takeaway(s, "Student is a class · Francis is an object of the Student class")

    # ---------------- 14. REAL WORLD -> OBJECTS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Business Concepts Become Objects", "The bank's real-world entities modelled as classes")
    _pic(s, "bank_objects.png", y=1.6, w=10.4)
    _takeaway(s, "Objects need not be physical — Registration, Payment and Order are objects too.")

    # ---------------- 15. ENCAPSULATION ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Encapsulation", "Protect data behind controlled methods")
    _pic(s, "encapsulation.png", y=1.65, w=10.6)
    _takeaway(s, "Never let the whole app write balance = -500000 directly.")

    # ---------------- 16. INHERITANCE & POLYMORPHISM ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Inheritance & Polymorphism", "Reuse common behaviour · respond differently to the same call")
    _pic(s, "inheritance.png", x=0.55, y=1.6, w=5.9)
    _pic(s, "polymorphism.png", x=6.7, y=1.7, w=6.1)
    _takeaway(s, "User → Student / Lecturer · displayDashboard() behaves differently per object")

    # ---------------- 17. FOUR PRINCIPLES ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "The Four OO Principles", "The ideas you will meet again and again")
    _pic(s, "four_principles.png", y=1.6, w=10.4)
    _takeaway(s, "Abstraction · Encapsulation · Inheritance · Polymorphism")

    # ---------------- 18. STRENGTHS & LIMITATIONS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Strengths & Limitations", "Object orientation is powerful — but it is not magic")
    _bullets(s, 0.9, 1.55, 5.9, 4.9, [
        ("Strengths: ", ""),
        "Modularity — systems divided into components",
        "Reusability — components and classes reused",
        "Maintainability — changes isolated",
        "Encapsulation — details hidden",
        "Extensibility — easy to grow",
        "Natural mapping of business concepts",
    ], size=15, gap=8)
    _bullets(s, 7.1, 1.55, 5.3, 4.9, [
        ("Limitations: ", ""),
        "Complexity — thousands of classes",
        "Learning curve — many concepts",
        "Poor modelling adds complexity",
        "Inheritance can be misused",
        "Up-front modelling effort",
    ], size=15, gap=10)
    _takeaway(s, "OOA → OOD → OOP are three different stages — analysis is not coding.")

    # ---------------- 19. ATM EXAMPLE ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Worked Example — the ATM", "From a business problem to candidate objects")
    _pic(s, "atm_flow.png", y=1.75, w=11.0)
    _takeaway(s, "business problem → business need → requirements → objects")

    # ---------------- 20. ACTIVITIES ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Practical & Tutorial", "Apply Week 1 to a manual library system")
    _bullets(s, 0.9, 1.55, 6.0, 4.9, [
        ("Practical (individual): ", "the college's manual library system."),
        "List ≥ 5 business problems", "List ≥ 3 business needs",
        "List ≥ 6 stakeholders", "List ≥ 8 functional requirements",
        "List ≥ 5 non-functional requirements", "List ≥ 8 candidate objects — justify each",
    ], size=16, gap=9)
    _bullets(s, 7.1, 1.55, 5.3, 4.9, [
        ("Tutorial (groups): ", "pick one system —"),
        "Mobile money", "Hospital management", "Online shopping",
        "Hotel booking", "University registration",
        "Produce: problem, need, objective, stakeholders, functional & non-functional needs, rules, objects.",
    ], size=16, gap=9)
    _takeaway(s, "Present your group's findings to the class.")

    # ---------------- 21. SUMMARY ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, CREAM)
    _header(s, "Week 1 Summary", "The definitions to remember")
    _bullets(s, 0.9, 1.55, 11.6, 4.9, [
        ("Business need ", "— a problem/opportunity/improvement the organisation must address."),
        ("Stakeholder ", "— anyone interested in or affected by the system."),
        ("Functional requirement ", "— what the system must do."),
        ("Non-functional requirement ", "— quality and constraints."),
        ("Object ", "— identity + state + behaviour."),
        ("Class ", "— the blueprint; object — the instance."),
    ], size=16, gap=12)
    _takeaway(s, "Business objective → problem → need → requirements → models → design → implementation.")

    # ---------------- 22. REVISION / THANKS ----------------
    s = prs.slides.add_slide(blank)
    _set_bg(s, MAROON)
    _rect(s, 0.35, 0.35, SW - 0.7, SH - 0.7, MAROON, line=GOLD)
    _text(s, 0.7, 1.3, SW - 1.4, 0.6, "REVISION & NEXT STEPS", size=30, color=GOLD,
          font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
    _bullets_col = _text(s, 0.9, 2.3, 11.5, 3.4, [
        "Attempt the revision questions in the Week 1 module and the assignment.",
        "Observe a system you use daily — write its problem, needs, stakeholders, requirements and objects.",
        "Identify 10 real-world objects around you; list their state and behaviour.",
        "Next week: we begin turning these ideas into formal analysis models.",
    ], size=18, color=CREAM)
    _text(s, 0.7, 6.3, SW - 1.4, 0.5, "Thank you — Francis Fweta · ICT Lecturer · Exploits University",
          size=14, color=GOLD_L, font=B_FONT, align=PP_ALIGN.CENTER)

    out = os.path.join(OUTPUT, "week1")
    os.makedirs(out, exist_ok=True)
    path = os.path.join(out, "Week1_Presentation_BICT3202_OOAD.pptx")
    prs.save(path)
    return path


if __name__ == "__main__":
    print("Built:", build())
