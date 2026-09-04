"""
PPTX engine — brand-styled presentation builder, parameterized by week.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from common import BRAND, VISUALS, OUTPUT

MAROON = RGBColor(0x7B, 0x11, 0x13)
DARK   = RGBColor(0x59, 0x09, 0x0C)
SOFT   = RGBColor(0xA9, 0x44, 0x42)
GOLD   = RGBColor(0xC9, 0xA2, 0x27)
GOLD_L = RGBColor(0xE4, 0xC9, 0x7C)
CREAM  = RGBColor(0xFA, 0xF6, 0xEC)
CREAM2 = RGBColor(0xF3, 0xEC, 0xDC)
TINT   = RGBColor(0xF3, 0xE4, 0xE4)
INK    = RGBColor(0x2B, 0x26, 0x26)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

H_FONT = "Poppins"
B_FONT = "Open Sans"

SW = 13.333
SH = 7.5


class Deck:
    def __init__(self, cfg):
        self.cfg = cfg
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.blank = self.prs.slide_layouts[6]
        self.seal = os.path.join(BRAND, "eu_seal.png")
        self._page = 0

    # ---- low level ----
    def bg(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def rect(self, slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(1)
        sp.shadow.inherit = False
        return sp

    def text(self, slide, x, y, w, h, runs, size=18, color=INK, font=B_FONT, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=0):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        if isinstance(runs, str):
            runs = [runs]
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.line_spacing = line_spacing
            if space_after:
                p.space_after = Pt(space_after)
            if isinstance(para, str):
                para = [(para, {})]
            for txt, ov in para:
                r = p.add_run()
                r.text = txt
                f = r.font
                f.name = ov.get("font", font)
                f.size = Pt(ov.get("size", size))
                f.bold = ov.get("bold", bold)
                f.italic = ov.get("italic", italic)
                f.color.rgb = ov.get("color", color)
        return box

    def pic(self, slide, name, x=None, y=None, w=None):
        path = os.path.join(VISUALS, name)
        if x is None:
            x = (SW - w) / 2.0
        return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

    def bullets(self, slide, x, y, w, h, items, size=16, gap=8, color=INK):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.line_spacing = 1.06
            p.space_after = Pt(gap)
            r = p.add_run()
            r.text = "▪  "
            r.font.name = B_FONT
            r.font.size = Pt(size)
            r.font.color.rgb = MAROON
            r.font.bold = True
            if isinstance(it, tuple):
                lead, rest = it
                r1 = p.add_run(); r1.text = lead
                r1.font.name = B_FONT; r1.font.size = Pt(size)
                r1.font.bold = True; r1.font.color.rgb = MAROON
                r2 = p.add_run(); r2.text = rest
                r2.font.name = B_FONT; r2.font.size = Pt(size)
                r2.font.color.rgb = color
            else:
                r1 = p.add_run(); r1.text = it
                r1.font.name = B_FONT; r1.font.size = Pt(size)
                r1.font.color.rgb = color
        return box

    # ---- high level ----
    def slide(self, bgcolor=CREAM):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, bgcolor)
        return s

    def header(self, s, title, subtitle=None):
        self._page += 1
        self.rect(s, 0, 0, SW, 1.18, MAROON)
        self.rect(s, 0, 1.18, SW, 0.06, GOLD)
        self.rect(s, 0, 1.24, SW, 0.02, DARK)
        self.text(s, 0.55, 0.16, 10.6, 0.6, title, size=27, color=WHITE,
                  font=H_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if subtitle:
            self.text(s, 0.55, 0.72, 10.6, 0.4, subtitle, size=13, color=GOLD_L,
                      font=B_FONT, anchor=MSO_ANCHOR.MIDDLE)
        self.text(s, 0.55, 7.12, 8, 0.3,
                  "BICT 3202 · Object-Oriented Analysis and Design", size=10, color=GREY)
        self.text(s, 9.5, 7.12, 3.3, 0.3,
                  f"Week {self.cfg['week']} · Slide {self._page}", size=10,
                  color=GREY, align=PP_ALIGN.RIGHT)

    def takeaway(self, s, text):
        self.rect(s, 0.9, 6.55, SW - 1.8, 0.5, DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self.text(s, 1.2, 6.55, SW - 2.4, 0.5, text, size=14, color=GOLD_L,
                  font=B_FONT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    def cover(self):
        s = self.slide(MAROON)
        self.rect(s, 0, 0, SW, 7.5, MAROON)
        self.rect(s, 0.35, 0.35, SW - 0.7, 6.8, MAROON, line=GOLD)
        self.rect(s, 0, 6.1, SW, 1.4, DARK)
        s.shapes.add_picture(self.seal, Inches((SW - 1.7) / 2), Inches(0.75), Inches(1.7), Inches(1.7))
        self.text(s, 0.5, 2.6, SW - 1, 0.6, "EXPLOITS UNIVERSITY", size=34, color=GOLD_L,
                  font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
        self.text(s, 0.5, 3.3, SW - 1, 0.4, "BSc INFORMATION COMMUNICATION AND TECHNOLOGY",
                  size=14, color=CREAM, font=B_FONT, align=PP_ALIGN.CENTER)
        self.text(s, 0.5, 3.85, SW - 1, 0.4, "BICT 3202 · OBJECT-ORIENTED ANALYSIS AND DESIGN",
                  size=16, color=GOLD, font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
        self.text(s, 0.5, 4.5, SW - 1, 1.0, f"WEEK {self.cfg['week']}", size=54, color=WHITE,
                  font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
        self.text(s, 0.5, 5.35, SW - 1, 0.5, self.cfg["title"], size=17, color=CREAM,
                  font=H_FONT, align=PP_ALIGN.CENTER)
        self.text(s, 0.5, 6.35, SW - 1, 0.4,
                  f"Francis Fweta  ·  ICT Lecturer  ·  Week {self.cfg['week']} of 16",
                  size=13, color=GOLD_L, font=B_FONT, align=PP_ALIGN.CENTER)
        return s

    def closing(self, heading="REVISION & NEXT STEPS", points=None, thanks="Thank you — Francis Fweta · ICT Lecturer · Exploits University"):
        s = self.slide(MAROON)
        self.rect(s, 0.35, 0.35, SW - 0.7, SH - 0.7, MAROON, line=GOLD)
        self.text(s, 0.7, 1.3, SW - 1.4, 0.6, heading, size=30, color=GOLD,
                  font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
        if points:
            self.bullets(s, 0.9, 2.3, 11.5, 3.4, points, size=18, gap=10, color=CREAM)
        self.text(s, 0.7, 6.3, SW - 1.4, 0.5, thanks, size=14, color=GOLD_L,
                  font=B_FONT, align=PP_ALIGN.CENTER)
        return s

    def save(self, folder, filename):
        out = os.path.join(OUTPUT, folder)
        os.makedirs(out, exist_ok=True)
        path = os.path.join(out, filename)
        self.prs.save(path)
        return path


def build(cfg, slides_builder):
    d = Deck(cfg)
    d.cover()
    slides_builder(d)
    return d.save(f"week{cfg['week']}", cfg["filename"])
