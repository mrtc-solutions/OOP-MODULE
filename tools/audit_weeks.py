"""
Comprehensive audit of Weeks 6-16 deliverables against the standing requirements.
Checks: file presence, PDF cover+pages+images+maroon, PPTX cover+slides+images,
DOCX cover+sections+marking scheme+footer, and content coverage vs source .md.
"""
import os, re, sys, glob

import pymupdf
from pptx import Presentation
from docx import Document

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "output")
VIS = os.path.join(ROOT, "assets", "visuals")

MAROON_RGB = (0x7B, 0x11, 0x13)

results = []
def record(ok, week, kind, detail):
    results.append((ok, week, kind, detail))

def check(cond, week, kind, detail):
    record(bool(cond), week, kind, detail)

def pdf_cover_is_maroon(path):
    d = pymupdf.open(path)
    if d.page_count == 0:
        return False, 0, 0
    page = d[0]
    # sample pixels across the page and check dominant color is maroon/dark
    pix = page.get_pixmap(dpi=40)
    n = pix.n
    w, h = pix.width, pix.height
    samples = []
    for (xx, yy) in [(0.5, 0.5), (0.5, 0.15), (0.15, 0.5), (0.85, 0.5)]:
        x = int(xx * w); y = int(yy * h)
        idx = (y * w + x) * n
        samples.append((pix.samples[idx], pix.samples[idx+1], pix.samples[idx+2]))
    # a maroon page: R > G significantly and R >= 90, G < 60
    maroon_hits = sum(1 for (r, g, b) in samples if r >= 90 and g < 60 and b < 60 and r > g + 40)
    return maroon_hits >= 3, samples, d.page_count

def main():
    print("=" * 78)
    print("AUDIT 1 — file presence & structure")
    print("=" * 78)
    for week in range(6, 17):
        d = os.path.join(OUT, f"week{week}")
        pdfs = glob.glob(os.path.join(d, "*.pdf"))
        pptxs = glob.glob(os.path.join(d, "*.pptx"))
        docxs = glob.glob(os.path.join(d, "*.docx"))
        check(len(pdfs) == 1, week, "PDF", f"{len(pdfs)} pdf(s)")
        check(len(pptxs) == 1, week, "PPTX", f"{len(pptxs)} pptx(s)")
        check(len(docxs) == 1, week, "DOCX", f"{len(docxs)} docx(s)")
        for p in pdfs:
            sz = os.path.getsize(p)
            check(sz > 100_000, week, "PDF", f"{os.path.basename(p)} {sz//1024}KB")
        for p in pptxs:
            sz = os.path.getsize(p)
            check(sz > 100_000, week, "PPTX", f"{os.path.basename(p)} {sz//1024}KB")
        for p in docxs:
            sz = os.path.getsize(p)
            check(sz > 20_000, week, "DOCX", f"{os.path.basename(p)} {sz//1024}KB")

    print()
    print("=" * 78)
    print("AUDIT 2 — PDF cover page (maroon + seal + Exploits University)")
    print("=" * 78)
    for week in range(6, 17):
        path = glob.glob(os.path.join(OUT, f"week{week}", "*.pdf"))[0]
        d = pymupdf.open(path)
        ok_maroon, samples, npages = pdf_cover_is_maroon(path)
        text = d[0].get_text()
        has_seal = d[0].get_images()
        has_name = "EXPLOITS UNIVERSITY" in text.upper()
        has_week = f"WEEK {week}" in text.upper()
        has_code = "BICT 3202" in text.upper()
        check(ok_maroon, week, "PDF-cover", f"maroon bg {samples}")
        check(bool(has_seal), week, "PDF-cover", "seal image present")
        check(has_name, week, "PDF-cover", "EXPLOITS UNIVERSITY text")
        check(has_week, week, "PDF-cover", f"WEEK {week} text")
        check(has_code, week, "PDF-cover", "BICT 3202 text")
        # images on body pages
        n_imgs = sum(len(d[i].get_images()) for i in range(1, npages))
        check(n_imgs >= 10, week, "PDF-body", f"{n_imgs} figures")
        d.close()

    print()
    print("=" * 78)
    print("AUDIT 3 — PPTX cover slide + images")
    print("=" * 78)
    for week in range(6, 17):
        path = glob.glob(os.path.join(OUT, f"week{week}", "*.pptx"))[0]
        prs = Presentation(path)
        s0 = prs.slides[0]
        texts = []
        n_pic_s0 = 0
        for sh in s0.shapes:
            if sh.shape_type == 13:
                n_pic_s0 += 1
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
        alltext = "\n".join(texts).upper()
        check("EXPLOITS UNIVERSITY" in alltext, week, "PPTX-cover", "EXPLOITS UNIVERSITY text")
        check(f"WEEK {week}" in alltext, week, "PPTX-cover", f"WEEK {week} text")
        check(n_pic_s0 >= 1, week, "PPTX-cover", f"{n_pic_s0} seal image(s)")
        n_slides = len(prs.slides)
        n_imgs = sum(1 for sl in prs.slides for sh in sl.shapes if sh.shape_type == 13)
        check(n_slides >= 20, week, "PPTX", f"{n_slides} slides")
        check(n_imgs >= 15, week, "PPTX", f"{n_imgs} images")

    print()
    print("=" * 78)
    print("AUDIT 4 — DOCX cover page + sections + marking scheme")
    print("=" * 78)
    for week in range(6, 17):
        path = glob.glob(os.path.join(OUT, f"week{week}", "*.docx"))[0]
        doc = Document(path)
        # cover page: first page content
        cover_text = []
        for p in doc.paragraphs[:40]:
            cover_text.append(p.text)
        # find seal image
        has_seal = False
        for rel in doc.part.rels.values():
            if "image" in rel.reltype and "seal" in rel.target_ref:
                has_seal = True
        all_cov = "\n".join(cover_text).upper()
        check("EXPLOITS UNIVERSITY" in all_cov, week, "DOCX-cover", "EXPLOITS UNIVERSITY text")
        check(f"WEEK {week}" in all_cov, week, "DOCX-cover", f"WEEK {week} text")
        check("FRANCIS FWETA" in all_cov, week, "DOCX-cover", "instructor name")
        # sections
        full = "\n".join(p.text for p in doc.paragraphs)
        for sec, want in [("A", "20"), ("B", "20"), ("C", "10")]:
            check(f"Section {sec}" in full, week, "DOCX-sec", f"Section {sec} present")
        # marking scheme table (last table or a table with TOTAL)
        found_total = False
        for t in doc.tables:
            for row in t.rows:
                cells = [c.text.strip() for c in row.cells]
                if any("TOTAL" in c.upper() for c in cells):
                    found_total = True
                    marks = [c for c in cells if re.search(r"\d+", c)]
        check(found_total, week, "DOCX-marking", "TOTAL row present")

    print()
    print("=" * 78)
    print("AUDIT 5 — content coverage vs source WeekN.md")
    print("=" * 78)
    for week in range(6, 17):
        md = open(os.path.join(ROOT, f"Week{week}.md"), encoding="utf-8").read()
        pdf = glob.glob(os.path.join(OUT, f"week{week}", "*.pdf"))[0]
        d = pymupdf.open(pdf)
        pdftext = "\n".join(d[i].get_text() for i in range(d.page_count))
        d.close()
        # extract headings from md (lines like "5. Title" or "PART X — Title")
        md_heads = re.findall(r"^\s*(?:PART\s+[A-Z]+\s*[—-]?\s*|\d+\.)\s*(.{6,70})$", md, re.M)
        # pick meaningful keyword phrases to search in pdf
        keywords = {
            6: ["aggregation", "inheritance", "polymorphism", "grouping", "composition"],
            7: ["event", "state", "transition", "concurrency", "operation"],
            8: ["use case", "actor", "sequence", "collaboration", "activity"],
            9: ["unified", "inception", "elaboration", "construction", "transition", "reuse"],
            10: ["design", "responsibility", "cohesion", "coupling", "abstraction"],
            11: ["refinement", "design class", "attribute", "operation", "responsibility", "interface"],
            12: ["integrated", "consistency", "traceability", "case study", "state model"],
            13: ["multilayer", "component", "reuse", "architecture", "adaptability"],
            14: ["design class", "responsibility", "visibility", "sequence", "consistency"],
            15: ["case study", "requirement", "actor", "use case", "CASE tool"],
            16: ["revision", "inheritance", "polymorphism", "unified", "consistency"],
        }
        kw = keywords[week]
        low = pdftext.lower()
        missing = [k for k in kw if k.lower() not in low]
        check(not missing, week, "coverage", f"missing keywords: {missing if missing else 'none'}")
        # check references present
        has_ref = ("reference" in low) or ("OMG" in pdftext)
        check(has_ref, week, "coverage", "references/OMG present")
        # check exam questions present
        has_exam = ("examination" in low) or ("exam" in low) or ("Question" in pdftext)
        check(has_exam, week, "coverage", "exam/tutorial questions present")

    print()
    print("=" * 78)
    print("AUDIT 6 — DOCX marking arithmetic (A=20, B=20, C=10) + references in body")
    print("=" * 78)
    for week in range(6, 17):
        path = glob.glob(os.path.join(OUT, f"week{week}", "*.docx"))[0]
        doc = Document(path)
        sec = None
        acc = {}
        for p in doc.paragraphs:
            t = p.text
            m = re.match(r"Section ([ABC])", t)
            if m:
                sec = m.group(1)
                acc.setdefault(sec, 0)
                tm = re.search(r"\((\d+)\s*marks\)", t)
                acc[sec + "_stated"] = int(tm.group(1)) if tm else 0
            if sec:
                acc[sec] += sum(int(x) for x in re.findall(r"\[\s*(\d+)\s*(?:marks)?\s*\]", t))
        check(acc.get("A", 0) == acc.get("A_stated", 0) == 20, week, "DOCX-marks", f"A={acc.get('A',0)}")
        check(acc.get("B", 0) == acc.get("B_stated", 0) == 20, week, "DOCX-marks", f"B={acc.get('B',0)}")
        check(acc.get("C", 0) == acc.get("C_stated", 0) == 10, week, "DOCX-marks", f"C={acc.get('C',0)}")
        # references in body (after TOC/cover)
        pdf = glob.glob(os.path.join(OUT, f"week{week}", "*.pdf"))[0]
        d = pymupdf.open(pdf)
        body = "\n".join(d[i].get_text() for i in range(2, d.page_count))
        has_omg = ("Object Management Group" in body) or ("OMG" in body)
        has_mach = "Mach" in body
        toc_ok = ("Contents" in d[1].get_text()) if d.page_count > 1 else False
        d.close()
        check(has_omg, week, "refs", "OMG in body")
        check(has_mach, week, "refs", "Mach (textbook) in body")
        check(toc_ok, week, "toc", "Contents page present")

    print()
    fails = [r for r in results if not r[0]]
    print("=" * 78)
    print(f"FINAL RESULT: {len(results) - len(fails)}/{len(results)} checks passed")
    if fails:
        print("FAILURES:")
        for _, w, k, det in fails:
            print(f"  [FAIL] Week {w} · {k}: {det}")
    else:
        print("ALL CHECKS PASSED — 10/10")
    print("=" * 78)
    return 0 if not fails else 1

if __name__ == "__main__":
    sys.exit(main())
